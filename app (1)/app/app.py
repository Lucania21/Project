from flask import Flask, render_template, request, redirect, url_for, make_response, jsonify
import re
from Backend import *
from docx import Document
from pptx import Presentation
import io
app = Flask(__name__)
@app.route('/')
def index():
    try:
        user = request.cookies.get('user')
        if not user:
            return redirect(url_for('login'))
    except:
        return redirect(url_for('login'))
    return render_template('index.html', user=user)

@app.route('/files')
def files():
    try:
        user = request.cookies.get('user')
        user_id = request.cookies.get('id')
        if not user or not user_id:
            return redirect(url_for('login'))
    except:
        return redirect(url_for('login'))

    files = DBretrieve('Filesmanager', ['id', 'name', 'date'], ID=['user_id', user])
    files= [{'user':user,'id':i[0],'name':i[1],'date':i[2]}for i in files]
    return render_template('files.html', files=files)

@app.route('/summaries')
def summaries():
    try:
        user = request.cookies.get('user')
        user_id = request.cookies.get('id')
        if not user or not user_id:
            return redirect(url_for('login'))
    except:
        return redirect(url_for('login'))

    summaries = DBretrieve('summaries', ['id', 'date'], ['userid', user, 'active', True])
    summaries=[{'user':user,'id':i[0],'name':'sub','date':i[1]}for i in summaries]
    return render_template('summaries.html', summaries=summaries)

@app.route('/quizzes')
def quizzes():
    try:
        user = request.cookies.get('user')
        user_id = request.cookies.get('id')
        if not user or not user_id:
            return redirect(url_for('login'))
    except:
        return redirect(url_for('login'))
    # print(['file_id', 'name', 'date_created', 'top_score', 'num_questions', 'type'],['userid', user_id, 'active', 'True'])
    quizzes = DBretrieve('Quizlet', ['file_id', 'name', 'date_created', 'top_score', 'num_questions', 'type'], ['userid', user_id, 'active', True])
    quiz_list = [{'id': i[0], 'name': i[1], 'created_date': i[2], 'num_questions': i[4], 'type': 'Multiple Choice'} for i in quizzes]
    
    return render_template('quizzes.html', quiz_list=quiz_list)

@app.route('/donate')
def donate():
    return render_template('donate.html')

@app.route('/about')
def about():
    return render_template('about.html')

@app.route('/contact')
def contact():
    return render_template('contact.html')

@app.route('/login', methods=['GET', 'POST'])
def login():
    error = None
    if request.method == 'POST':
        username = request.form['username']
        password = request.form['password']
        
        user_data = DBretrieve('User', ['id', 'passwordhash'], ID=['username', username])
        if not user_data:
            error = 'Username or password incorrect'
        else:
            user_id, truepass = user_data[0]
            if password != truepass:
                error = 'Username or password incorrect'
            else:
                resp = make_response(redirect(url_for('index')))
                resp.set_cookie('user', username)
                resp.set_cookie('id', str(user_id))
                return resp

    return render_template('login.html', error=error)

def is_valid_email(email):
    return re.match(r"[^@]+@[^@]+\.[^@]+", email)

@app.route('/signup', methods=['GET', 'POST'])
def signup():
    error = None
    success = None
    if request.method == 'POST':
        username = request.form['username']
        email = request.form['email']
        password = request.form['password']
        confirm_password = request.form['confirm_password']

        if not username or not email or not password or not confirm_password:
            error = 'All fields are required.'
        elif password != confirm_password:
            error = 'Passwords do not match.'
        elif not is_valid_email(email):
            error = 'Invalid email address.'
        elif len(password) < 8:
            error = 'Password must be at least 8 characters long.'
        else:
            try:
                DBupload('User', [None, username, email, password, True])
                success = 'Account created successfully! Please log in.'
                return redirect(url_for('login'))
            except Exception as e:
                error = 'An error occurred while creating your account. Please try again.'
                print(f"Error: {e}")

    return render_template('signup.html', error=error, success=success)

@app.route('/delete_summary/<summary_id>', methods=['GET'])
def delete_summary(summary_id):
    DBedit('summaries', {"active": False}, ['id', summary_id])
    return redirect(url_for('summaries'))
@app.route('/delete_quiz/<quiz>', methods=['DELETE'])
def delete_quiz(quiz):
    DBedit('Quizlet', {"active": False}, ['id', quiz])
    return redirect(url_for('quizzes'))
@app.route('/summary/<summary_id>')
def summary(summary_id):
    summary_path = DBretrieve('summaries', ['filepath'], ['id', summary_id])
    if summary_path:
        summary_path = summary_path[0][0][:-1][1:]
        with open(summary_path, 'r') as file:
            summary_content = file.read()
    else:
        summary_content = "Summary not found."
    return render_template('summary.html', summary_content=summary_content)

@app.route('/quiz/<quiz_id>')
def quiz(quiz_id):
    quiz_data = DBretrieve('Questions', ['id', 'Questiontxt', 'wrongoption', 'correctoption'], ['Quizlet_id', quiz_id])
    print(quiz_data)
    quiz = [{'question': q[1], 'answer': q[3], 'wrong_option': q[2].replace('[','').replace(']','').replace("'",'').split(',')} for q in quiz_data]
    quiz_name = DBretrieve('Quizlet', ['name'], ['file_id', quiz_id])[0][0]
    return render_template('quiz.html', quiz_data=quiz, quiz_name=quiz_name)

@app.route('/submit_quiz', methods=['POST'])
def submit_quiz():
    data = request.json
    score = data.get('score')
    quiz_id = data.get('quiz_id')
    DBedit('Quizlet', {"top_score": score}, ['file_id', quiz_id])
    return jsonify({'status': 'success', 'score': score})

@app.route('/delete_file/<int:file_id>', methods=['DELETE'])
def delete_file(file_id):
    DBedit('Filesmanager', {"active": False}, ['id', file_id])
    return jsonify({'status': 'success'})

@app.route('/create_quiz')
def create_quiz():
    user = request.cookies.get('user')
    user_id = request.cookies.get('id')
    if not user or not user_id:
        return redirect(url_for('login'))   
    file_id = request.args.get('file_id')
    Generatequiz(file_id,user_id)
    return redirect(url_for('quizzes'))


@app.route('/create_summary')
def create_summary():
    file_id = request.args.get('file_id')
    summary_type = request.args.get('type')
    if summary_type == '2':
        Generatesum(file_id, 1, request.cookies.get('user'))
        return redirect(url_for('summaries'))    
    elif summary_type == '1':
        Generatesum(file_id, 0, request.cookies.get('user'))
        return redirect(url_for('summaries'))
    else:
        return "Invalid summary type", 400

def extract_text_from_docx(file_stream):
    doc = Document(file_stream)
    full_text = []
    for paragraph in doc.paragraphs:
        full_text.append(paragraph.text)
    return '\n'.join(full_text)
def extract_text_from_pptx(file):
    prs = Presentation(io.BytesIO(file.read()))
    full_text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                full_text.append(shape.text)
    return '\n'.join(full_text)
@app.route('/upload_new_file', methods=['POST'])
def upload_new_file():
    if 'file' not in request.files:
        return jsonify({'error': 'No file part'}), 400
    
    file = request.files['file']
    file_name = file.filename
    if file.filename == '':
        return jsonify({'error': 'No selected file'}), 400

    try:
        if file.filename.endswith('.txt'):
            file_content = file.read().decode('utf-8')
        elif file.filename.endswith('.docx'):
            file_content = extract_text_from_docx(file)
        elif file_name.endswith('.pptx'):
            file_content = extract_text_from_pptx(file)        
        else:
            return jsonify({'error': 'Unsupported file type'}), 400
        user = request.cookies.get('user')
        category = request.form.get('category')
        Uploadfile(user, category,file_name, file_content)

        return jsonify({'status': 'success'})
    except Exception as e:
        return jsonify({'error': str(e)}), 500

if __name__ == '__main__':
    app.run(debug=True)
